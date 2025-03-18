import logging
import os
from typing import Any

from ....core.exceptions import ExtractorError
from ..extractor import Extractor

logger = logging.getLogger(__name__)


class DocumentExtractor(Extractor):
    """Extractor for document files (DOC, DOCX, XLS, XLSX, PPT, PPTX, etc.)."""

    def __init__(self):
        super().__init__()
        self._supported_extensions = {
            # Microsoft Office
            ".doc",
            ".docx",
            ".xls",
            ".xlsx",
            ".ppt",
            ".pptx",
            # OpenOffice / LibreOffice
            ".odt",
            ".ods",
            ".odp",
            # Other document formats
            ".rtf",
            ".epub",
        }
        self._supported_mime_types = {
            # Microsoft Office
            "application/msword",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/vnd.ms-excel",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-powerpoint",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            # OpenOffice / LibreOffice
            "application/vnd.oasis.opendocument.text",
            "application/vnd.oasis.opendocument.spreadsheet",
            "application/vnd.oasis.opendocument.presentation",
            # Other
            "application/rtf",
            "application/epub+zip",
        }

        # Check if required libraries are available
        self._has_python_docx = self._check_python_docx()
        self._has_xlrd = self._check_xlrd()
        self._has_python_pptx = self._check_python_pptx()
        self._has_odf = self._check_odf()
        self._has_ebooklib = self._check_ebooklib()

        if not any(
            [
                self._has_python_docx,
                self._has_xlrd,
                self._has_python_pptx,
                self._has_odf,
                self._has_ebooklib,
            ]
        ):
            logger.warning(
                "No document extraction libraries found. Document metadata extraction will be limited."
            )

    def _check_python_docx(self) -> bool:
        try:
            import docx

            return True
        except ImportError:
            return False

    def _check_xlrd(self) -> bool:
        try:
            import xlrd

            return True
        except ImportError:
            return False

    def _check_python_pptx(self) -> bool:
        try:
            import pptx

            return True
        except ImportError:
            return False

    def _check_odf(self) -> bool:
        try:
            import odf

            return True
        except ImportError:
            return False

    def _check_ebooklib(self) -> bool:
        try:
            import ebooklib

            return True
        except ImportError:
            return False

    def extract_metadata(self, file_path: str) -> dict[str, Any]:
        """Extract metadata from a document file.

        Args:
            file_path: Path to the document file

        Returns:
            Dictionary containing metadata

        Raises:
            ExtractorError: If metadata extraction fails
        """
        if not os.path.exists(file_path):
            raise ExtractorError(f"File not found: {file_path}")

        try:
            # Basic file info
            stat_info = os.stat(file_path)
            _, ext = os.path.splitext(file_path)
            ext_lower = ext.lower()

            metadata = {
                "file_type": "document",
                "extension": ext_lower,
                "size_bytes": stat_info.st_size,
                "last_modified": stat_info.st_mtime,
                "created": stat_info.st_ctime,
            }

            # Enhanced metadata based on file type
            if ext_lower in [".docx"] and self._has_python_docx:
                metadata.update(self._extract_docx_metadata(file_path))
            elif ext_lower in [".xlsx", ".xls"] and self._has_xlrd:
                metadata.update(self._extract_excel_metadata(file_path))
            elif ext_lower in [".pptx"] and self._has_python_pptx:
                metadata.update(self._extract_pptx_metadata(file_path))
            elif ext_lower in [".odt", ".ods", ".odp"] and self._has_odf:
                metadata.update(self._extract_odf_metadata(file_path))
            elif ext_lower in [".epub"] and self._has_ebooklib:
                metadata.update(self._extract_epub_metadata(file_path))

            return metadata

        except Exception as e:
            logger.error(f"Error extracting metadata from {file_path}: {str(e)}")
            raise ExtractorError(
                f"Failed to extract metadata from {file_path}: {str(e)}"
            ) from e

    def extract_content(self, file_path: str) -> str:
        """Extract textual content from a document file.

        Args:
            file_path: Path to the document file

        Returns:
            Extracted textual content

        Raises:
            ExtractorError: If content extraction fails
        """
        if not os.path.exists(file_path):
            raise ExtractorError(f"File not found: {file_path}")

        try:
            _, ext = os.path.splitext(file_path)
            ext_lower = ext.lower()

            if ext_lower in [".docx"] and self._has_python_docx:
                return self._extract_docx_content(file_path)
            elif ext_lower in [".xlsx", ".xls"] and self._has_xlrd:
                return self._extract_excel_content(file_path)
            elif ext_lower in [".pptx"] and self._has_python_pptx:
                return self._extract_pptx_content(file_path)
            elif ext_lower in [".odt", ".ods", ".odp"] and self._has_odf:
                return self._extract_odf_content(file_path)
            elif ext_lower in [".epub"] and self._has_ebooklib:
                return self._extract_epub_content(file_path)
            else:
                return f"[Document: {os.path.basename(file_path)}]\nContent extraction not supported for this file type."

        except Exception as e:
            logger.error(f"Error extracting content from {file_path}: {str(e)}")
            raise ExtractorError(
                f"Failed to extract content from {file_path}: {str(e)}"
            ) from e

    def generate_preview(self, file_path: str, max_size: int = 1000) -> str:
        """Generate a textual preview of the document file.

        Args:
            file_path: Path to the document file
            max_size: Maximum size of the preview in characters

        Returns:
            A textual preview of the document content

        Raises:
            ExtractorError: If preview generation fails
        """
        try:
            # Get metadata
            metadata = self.extract_metadata(file_path)

            preview = f"[Document: {os.path.basename(file_path)}]\n"

            # Add basic metadata
            if "title" in metadata:
                preview += f"Title: {metadata['title']}\n"
            if "author" in metadata:
                preview += f"Author: {metadata['author']}\n"
            if "created_date" in metadata:
                preview += f"Created: {metadata['created_date']}\n"
            if "page_count" in metadata:
                preview += f"Pages: {metadata['page_count']}\n"
            if "word_count" in metadata:
                preview += f"Words: {metadata['word_count']}\n"

            preview += "\nContent:\n"

            # Add extracted content
            content = self.extract_content(file_path)
            remaining_size = max_size - len(preview)

            if remaining_size <= 0:
                return preview

            if len(content) <= remaining_size:
                preview += content
            else:
                preview += content[:remaining_size] + "..."

            return preview

        except Exception as e:
            logger.error(f"Error generating preview for {file_path}: {str(e)}")
            raise ExtractorError(
                f"Failed to generate preview for {file_path}: {str(e)}"
            ) from e

    def _extract_docx_metadata(self, file_path: str) -> dict[str, Any]:
        """Extract metadata from a DOCX file.

        Args:
            file_path: Path to the DOCX file

        Returns:
            Dictionary containing metadata
        """
        import docx

        metadata = {}
        doc = docx.Document(file_path)

        # Document properties
        core_properties = doc.core_properties
        if core_properties.title:
            metadata["title"] = core_properties.title
        if core_properties.author:
            metadata["author"] = core_properties.author
        if core_properties.created:
            metadata["created_date"] = core_properties.created
        if core_properties.modified:
            metadata["modified_date"] = core_properties.modified
        if core_properties.subject:
            metadata["subject"] = core_properties.subject
        if core_properties.keywords:
            metadata["keywords"] = core_properties.keywords
        if core_properties.category:
            metadata["category"] = core_properties.category
        if core_properties.comments:
            metadata["comments"] = core_properties.comments

        # Document statistics
        metadata["paragraph_count"] = len(doc.paragraphs)
        metadata["table_count"] = len(doc.tables)

        # Count words
        word_count = 0
        for para in doc.paragraphs:
            word_count += len(para.text.split())
        metadata["word_count"] = word_count

        return metadata

    def _extract_docx_content(self, file_path: str) -> str:
        """Extract content from a DOCX file.

        Args:
            file_path: Path to the DOCX file

        Returns:
            Extracted text content
        """
        import docx

        doc = docx.Document(file_path)

        # Extract text from paragraphs
        content = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                content.append(paragraph.text)

        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    content.append(" | ".join(row_text))

        return "\n\n".join(content)

    def _extract_excel_metadata(self, file_path: str) -> dict[str, Any]:
        """Extract metadata from an Excel file.

        Args:
            file_path: Path to the Excel file

        Returns:
            Dictionary containing metadata
        """
        import xlrd

        metadata = {}
        workbook = xlrd.open_workbook(file_path, on_demand=True)

        # Workbook properties
        metadata["sheet_count"] = workbook.nsheets
        metadata["sheet_names"] = workbook.sheet_names()

        # Count cells with data
        cell_count = 0
        for sheet_name in workbook.sheet_names():
            sheet = workbook.sheet_by_name(sheet_name)
            for row in range(sheet.nrows):
                for col in range(sheet.ncols):
                    if sheet.cell_value(row, col):
                        cell_count += 1

        metadata["cell_count"] = cell_count

        return metadata

    def _extract_excel_content(self, file_path: str) -> str:
        """Extract content from an Excel file.

        Args:
            file_path: Path to the Excel file

        Returns:
            Extracted text content
        """
        import xlrd

        workbook = xlrd.open_workbook(file_path)
        content = []

        for sheet_name in workbook.sheet_names():
            sheet = workbook.sheet_by_name(sheet_name)
            content.append(f"Sheet: {sheet_name}")

            # Extract data from cells
            for row in range(
                min(sheet.nrows, 100)
            ):  # Limit to 100 rows for performance
                row_values = []
                for col in range(sheet.ncols):
                    cell_value = sheet.cell_value(row, col)
                    if cell_value:
                        row_values.append(str(cell_value))
                if row_values:
                    content.append(" | ".join(row_values))

            if sheet.nrows > 100:
                content.append("... (more rows not shown)")

            content.append("")

        return "\n".join(content)

    def _extract_pptx_metadata(self, file_path: str) -> dict[str, Any]:
        """Extract metadata from a PPTX file.

        Args:
            file_path: Path to the PPTX file

        Returns:
            Dictionary containing metadata
        """
        import pptx

        metadata = {}
        presentation = pptx.Presentation(file_path)

        # Presentation properties
        if (
            hasattr(presentation.core_properties, "title")
            and presentation.core_properties.title
        ):
            metadata["title"] = presentation.core_properties.title
        if (
            hasattr(presentation.core_properties, "author")
            and presentation.core_properties.author
        ):
            metadata["author"] = presentation.core_properties.author
        if (
            hasattr(presentation.core_properties, "created")
            and presentation.core_properties.created
        ):
            metadata["created_date"] = presentation.core_properties.created
        if (
            hasattr(presentation.core_properties, "modified")
            and presentation.core_properties.modified
        ):
            metadata["modified_date"] = presentation.core_properties.modified

        # Presentation statistics
        metadata["slide_count"] = len(presentation.slides)

        return metadata

    def _extract_pptx_content(self, file_path: str) -> str:
        """Extract content from a PPTX file.

        Args:
            file_path: Path to the PPTX file

        Returns:
            Extracted text content
        """
        import pptx

        presentation = pptx.Presentation(file_path)
        content = []

        for i, slide in enumerate(presentation.slides):
            content.append(f"Slide {i + 1}:")

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content.append(shape.text.strip())

            content.append("")

        return "\n".join(content)

    def _extract_odf_metadata(self, file_path: str) -> dict[str, Any]:
        """Extract metadata from an ODF file (ODT, ODS, ODP).

        Args:
            file_path: Path to the ODF file

        Returns:
            Dictionary containing metadata
        """
        from odf import meta, opendocument

        metadata = {}
        doc = opendocument.load(file_path)

        # Document metadata
        doc_meta = meta.Meta(doc)
        if doc_meta.title:
            metadata["title"] = doc_meta.title
        if doc_meta.creator:
            metadata["author"] = doc_meta.creator
        if doc_meta.creation_date:
            metadata["created_date"] = doc_meta.creation_date
        if doc_meta.modification_date:
            metadata["modified_date"] = doc_meta.modification_date
        if doc_meta.subject:
            metadata["subject"] = doc_meta.subject
        if doc_meta.keywords:
            metadata["keywords"] = doc_meta.keywords
        if doc_meta.description:
            metadata["description"] = doc_meta.description

        return metadata

    def _extract_odf_content(self, file_path: str) -> str:
        """Extract content from an ODF file (ODT, ODS, ODP).

        Args:
            file_path: Path to the ODF file

        Returns:
            Extracted text content
        """
        from odf import opendocument, text

        doc = opendocument.load(file_path)
        content = []

        # Extract text using plain text export
        all_text = text.Plain.extract(doc.body)
        content.append(all_text)

        return "\n".join(content)

    def _extract_epub_metadata(self, file_path: str) -> dict[str, Any]:
        """Extract metadata from an EPUB file.

        Args:
            file_path: Path to the EPUB file

        Returns:
            Dictionary containing metadata
        """
        from ebooklib import epub

        metadata = {}
        book = epub.read_epub(file_path)

        # Book metadata
        if book.title:
            metadata["title"] = book.title
        if book.get_metadata("DC", "creator"):
            metadata["author"] = book.get_metadata("DC", "creator")[0][0]
        if book.get_metadata("DC", "publisher"):
            metadata["publisher"] = book.get_metadata("DC", "publisher")[0][0]
        if book.get_metadata("DC", "date"):
            metadata["published_date"] = book.get_metadata("DC", "date")[0][0]
        if book.get_metadata("DC", "language"):
            metadata["language"] = book.get_metadata("DC", "language")[0][0]
        if book.get_metadata("DC", "identifier"):
            for identifier in book.get_metadata("DC", "identifier"):
                if identifier[1].get("id") == "isbn":
                    metadata["isbn"] = identifier[0]

        # Count items
        metadata["item_count"] = len(book.get_items())

        return metadata

    def _extract_epub_content(self, file_path: str) -> str:
        """Extract content from an EPUB file.

        Args:
            file_path: Path to the EPUB file

        Returns:
            Extracted text content
        """
        import re

        import ebooklib
        from ebooklib import epub

        book = epub.read_epub(file_path)
        content = []

        # Add title
        if book.title:
            content.append(f"Title: {book.title}")

        # Extract content from HTML documents
        try:
            chapters = []
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    # Get HTML content
                    html_content = item.get_content().decode("utf-8")

                    # Strip HTML tags
                    text = re.sub("<[^<]+?>", " ", html_content)
                    text = re.sub(r"\s+", " ", text).strip()

                    if text:
                        chapters.append(text)

            # Add chapters
            for i, chapter in enumerate(chapters):
                content.append(f"\nChapter {i + 1}:")
                content.append(
                    chapter[:2000] + "..." if len(chapter) > 2000 else chapter
                )
        except Exception as e:
            content.append(f"Error extracting content: {str(e)}")

        return "\n".join(content)
